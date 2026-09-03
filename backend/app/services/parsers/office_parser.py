import io
import csv
from typing import Tuple, Dict, Any
from fastapi import HTTPException, status
from app.services.parsers.base import BaseDocumentParser


class OfficeDocumentParser(BaseDocumentParser):
    def parse(self, content: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

        if ext == "pptx":
            return self._parse_pptx(content)
        elif ext == "xlsx":
            return self._parse_xlsx(content)
        elif ext == "csv":
            return self._parse_csv(content)
        else:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Unsupported office extension '{ext}'."
            )

    def _parse_pptx(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Failed to parse PPTX presentation: {str(e)}"
            )

        slides_text = []
        for idx, slide in enumerate(prs.slides):
            slide_lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        txt = paragraph.text.strip()
                        if txt:
                            slide_lines.append(txt)
            if slide_lines:
                slides_text.append(f"## Slide {idx + 1}\n" + "\n".join(slide_lines))

        full_text = "\n\n".join(slides_text).strip()
        if not full_text:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="PPTX presentation contains no extractable text content."
            )

        return full_text, {
            "format": "pptx",
            "slide_count": len(prs.slides)
        }

    def _parse_xlsx(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Failed to parse XLSX workbook: {str(e)}"
            )

        sheets_text = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_rows = []
            for row in sheet.iter_rows(values_only=True):
                non_empty = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if non_empty:
                    sheet_rows.append(" | ".join(non_empty))
            if sheet_rows:
                sheets_text.append(f"## Sheet: {sheet_name}\n" + "\n".join(sheet_rows))

        full_text = "\n\n".join(sheets_text).strip()
        if not full_text:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="XLSX workbook contains no extractable data."
            )

        return full_text, {
            "format": "xlsx",
            "sheet_count": len(wb.sheetnames)
        }

    def _parse_csv(self, content: bytes) -> Tuple[str, Dict[str, Any]]:
        try:
            decoded = content.decode("utf-8")
        except UnicodeDecodeError:
            decoded = content.decode("latin-1")

        reader = csv.reader(io.StringIO(decoded))
        rows = [" | ".join(cell.strip() for cell in row if cell.strip()) for row in reader]
        rows = [r for r in rows if r]

        full_text = "\n".join(rows).strip()
        if not full_text:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="CSV document contains no extractable rows."
            )

        return full_text, {
            "format": "csv",
            "row_count": len(rows)
        }
